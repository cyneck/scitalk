"""PPT API"""
import os
import uuid
import shutil
from pathlib import Path

from fastapi import APIRouter, HTTPException, UploadFile, File
from pydantic import BaseModel

from app.core.config import get_settings

router = APIRouter(prefix="/ppt")

settings = get_settings()


class SlideResponse(BaseModel):
    """幻灯片响应"""
    id: str
    index: int
    image_path: str | None
    content_text: str


class PPTUploadResponse(BaseModel):
    """PPT上传响应"""
    project_id: str
    filename: str
    slides_count: int


@router.post("/upload", response_model=PPTUploadResponse)
async def upload_ppt(file: UploadFile = File(...)):
    """上传PPT文件"""
    if not file.filename:
        raise HTTPException(status_code=400, detail="未提供文件名")

    # 确保存储目录存在
    os.makedirs(settings.ppt_path, exist_ok=True)
    os.makedirs(settings.slides_path, exist_ok=True)

    # 保存文件
    project_id = str(uuid.uuid4())
    filename = f"{project_id}_{file.filename}"
    filepath = Path(settings.ppt_path) / filename

    with open(filepath, "wb") as f:
        content = await file.read()
        f.write(content)

    # 解析PPT（简化版，实际需要调用python-pptx）
    slides_count = await parse_ppt(filepath, project_id)

    return PPTUploadResponse(
        project_id=project_id,
        filename=filename,
        slides_count=slides_count,
    )


async def parse_ppt(pptx_path: Path, project_id: str) -> int:
    """解析PPT文件，提取幻灯片"""
    # TODO: 实现完整的PPT解析
    # 这里简化处理，实际需要：
    # 1. 使用python-pptx读取pptx文件
    # 2. 将每张幻灯片渲染为图片
    # 3. 提取文本内容

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
        from PIL import Image
        import io

        prs = Presentation(str(pptx_path))
        slides_info = []

        for i, slide in enumerate(prs.slides):
            # 提取文本
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)

            # 生成幻灯片图片
            # 简化处理：创建占位图片
            slide_id = str(uuid.uuid4())
            slide_image_path = Path(settings.slides_path) / f"{project_id}_{i}.png"

            # 创建空白图片作为占位
            img = Image.new('RGB', (1280, 720), color='white')
            img.save(slide_image_path)

            slides_info.append({
                "id": slide_id,
                "index": i,
                "image_path": str(slide_image_path),
                "content_text": "\n".join(text_content),
            })

        # 保存slides信息到临时存储
        _slides_storage[project_id] = slides_info

        return len(slides_info)

    except ImportError:
        # 如果没有python-pptx，返回模拟数据
        slides_info = []
        for i in range(5):
            slide_id = str(uuid.uuid4())
            slide_image_path = Path(settings.slides_path) / f"{project_id}_{i}.png"

            # 创建占位图片
            from PIL import Image
            img = Image.new('RGB', (1280, 720), color='lightgray')
            img.save(slide_image_path)

            slides_info.append({
                "id": slide_id,
                "index": i,
                "image_path": str(slide_image_path),
                "content_text": f"这是第 {i+1} 张幻灯片的内容",
            })

        _slides_storage[project_id] = slides_info
        return len(slides_info)


# 临时存储slides信息
_slides_storage: dict[str, list] = {}


@router.get("/{project_id}/slides")
async def get_slides(project_id: str):
    """获取项目的幻灯片列表"""
    if project_id not in _slides_storage:
        raise HTTPException(status_code=404, detail="项目不存在或PPT未上传")
    return _slides_storage[project_id]


@router.get("/{project_id}/slides/{slide_index}")
async def get_slide(project_id: str, slide_index: int):
    """获取指定幻灯片"""
    if project_id not in _slides_storage:
        raise HTTPException(status_code=404, detail="项目不存在")

    slides = _slides_storage[project_id]
    if slide_index >= len(slides):
        raise HTTPException(status_code=404, detail="幻灯片不存在")

    return slides[slide_index]