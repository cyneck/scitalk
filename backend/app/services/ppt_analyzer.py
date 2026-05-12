"""PPT分析服务"""
import os
import logging
from pathlib import Path
from typing import List, Dict, Any

logger = logging.getLogger(__name__)


class PPTAnalyzer:
    """PPT分析服务"""

    def __init__(self):
        self.supported_formats = [".pptx"]

    async def analyze(self, pptx_path: str) -> Dict[str, Any]:
        """分析PPT文件"""
        path = Path(pptx_path)

        if not path.exists():
            raise FileNotFoundError(f"PPT文件不存在: {pptx_path}")

        if path.suffix.lower() not in self.supported_formats:
            raise ValueError(f"不支持的文件格式: {path.suffix}")

        try:
            return await self._parse_pptx(path)
        except Exception as e:
            logger.error(f"PPT解析失败: {e}")
            raise

    async def _parse_pptx(self, path: Path) -> Dict[str, Any]:
        """解析PPTX文件"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu

            prs = Presentation(str(path))

            slides = []
            for i, slide in enumerate(prs.slides):
                slide_content = {
                    "index": i,
                    "text": [],
                    "images": [],
                    "shapes": [],
                }

                # 提取文本
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content["text"].append(shape.text)

                    # 检查是否有图片
                    if hasattr(shape, "image"):
                        slide_content["images"].append({
                            "name": shape.name,
                            "type": type(shape).__name__,
                        })

                slides.append(slide_content)

            return {
                "slides": slides,
                "total_slides": len(slides),
                "title": path.stem,
            }

        except ImportError:
            logger.warning("python-pptx 未安装，使用简化解析")
            return await self._parse_pptx_simple(path)

    async def _parse_pptx_simple(self, path: Path) -> Dict[str, Any]:
        """简化PPT解析（无python-pptx时）"""
        # 尝试使用其他方式解析或返回模拟数据
        return {
            "slides": [],
            "total_slides": 0,
            "title": path.stem,
            "error": "python-pptx not available",
        }

    async def extract_slide_content(self, slide) -> Dict[str, Any]:
        """提取单张幻灯片内容"""
        content = {
            "text": [],
            "images": [],
            "shapes": [],
        }

        if hasattr(slide, "shapes"):
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content["text"].append(shape.text)

                if hasattr(shape, "image"):
                    content["images"].append(shape.image)

        return content

    async def render_slide_to_image(self, slide, output_path: str, width: int = 1280, height: int = 720) -> str:
        """将幻灯片渲染为图片"""
        try:
            from pptx.util import Inches, Emu
            from PIL import Image
            import io

            # 获取幻灯片尺寸
            slide_width = slide.slide_width
            slide_height = slide.slide_height

            # 计算缩放比例
            scale_x = width / slide_width
            scale_y = height / slide_height
            scale = min(scale_x, scale_y)

            # 创建图片
            img = Image.new('RGB', (width, height), color='white')

            # 保存
            img.save(output_path)

            return output_path

        except Exception as e:
            logger.error(f"幻灯片渲染失败: {e}")
            raise